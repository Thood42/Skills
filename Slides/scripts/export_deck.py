#!/usr/bin/env python3
"""
export_deck.py — turn a rendered deck into a shareable PDF and/or PowerPoint.

The deck is HTML; people often need a file to email or drop into a meeting. This takes
the per-slide PNGs that render.sh produces and packages them, one slide per page:

  PDF   — one 16:9 page per slide (lossless; via img2pdf).
  PPTX  — one slide per PowerPoint slide, image full-bleed on a 13.333x7.5in canvas.
          (Slides are images, not editable text — it's for distribution, not re-editing.)

USAGE
  # 1) render first (needs Chrome), then export from the PNGs:
  ./render.sh mydeck.html 12 shots
  python3 scripts/export_deck.py --slides shots --pdf mydeck.pdf --pptx mydeck.pptx

  # or let export render for you:
  python3 scripts/export_deck.py --deck mydeck.html --slides 12 --pdf mydeck.pdf

NOTES
  - Use the slide PNGs (slide-01.png …), not the _contact.png / _golden artifacts; this
    script already filters to slide-*.png in order.
  - PDF needs `img2pdf` (preferred) or Pillow; PPTX needs `python-pptx`. Both are pip-installable.
"""
import argparse, re, subprocess, sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent


def slide_pngs(d):
    d = Path(d)
    files = sorted(d.glob("slide-*.png"), key=lambda p: int(re.search(r"(\d+)", p.stem).group(1)))
    if not files:
        sys.exit(f"No slide-*.png found in {d}. Run render.sh first.")
    return files


def do_render(deck, n, outdir):
    rs = ROOT / "render.sh"
    print(f"Rendering {deck} ({n} slides) → {outdir}/ …")
    subprocess.run(["bash", str(rs), str(deck), str(n), str(outdir)], check=True)


def to_pdf(pngs, out):
    out = Path(out)
    try:
        import img2pdf
        out.write_bytes(img2pdf.convert([str(p) for p in pngs]))
    except ImportError:
        from PIL import Image
        imgs = [Image.open(p).convert("RGB") for p in pngs]
        imgs[0].save(out, save_all=True, append_images=imgs[1:])
    print(f"PDF  → {out}  ({len(pngs)} pages)")


def to_pptx(pngs, out):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("PPTX export needs python-pptx (pip install python-pptx).")
    from PIL import Image
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height
    for p in pngs:
        slide = prs.slides.add_slide(blank)
        iw, ih = Image.open(p).size
        # cover-fit the 16:9 canvas (our slides are 1280x720 = exact 16:9, so this is 1:1)
        scale = max(SW / iw, SH / ih)
        w, h = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(str(p), int((SW - w) / 2), int((SH - h) / 2), w, h)
    prs.save(out)
    print(f"PPTX → {out}  ({len(pngs)} slides)")


def main():
    ap = argparse.ArgumentParser(description="Export a rendered deck to PDF / PPTX.")
    ap.add_argument("--slides", help="directory of slide-*.png (output of render.sh)")
    ap.add_argument("--deck", help="deck .html to render first (needs Chrome)")
    ap.add_argument("--n", "--count", dest="n", type=int, help="slide count (with --deck)")
    ap.add_argument("--pdf"); ap.add_argument("--pptx")
    a = ap.parse_args()

    if not a.pdf and not a.pptx:
        ap.error("nothing to do: pass --pdf and/or --pptx")

    shots = a.slides
    if a.deck:
        if not a.n:
            ap.error("--deck requires --n <slide count>")
        shots = a.slides or "shots"
        do_render(a.deck, a.n, shots)
    if not shots:
        ap.error("pass --slides <dir> (or --deck + --n to render first)")

    pngs = slide_pngs(shots)
    if a.pdf:
        to_pdf(pngs, a.pdf)
    if a.pptx:
        to_pptx(pngs, a.pptx)


if __name__ == "__main__":
    main()
